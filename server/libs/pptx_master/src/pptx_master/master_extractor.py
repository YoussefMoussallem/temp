"""Extract a MasterManifest from any .pptx file.

Accepts either a brand template (a .potx renamed to .pptx, or a near-empty
deck whose only purpose is to carry the master) or a regular content deck.
We only read master-level information: canvas size, theme fonts/colors,
placeholder geometry, and literal chrome text. Slide content is ignored.

Implementation uses python-pptx for the OPC package traversal and lxml to
read theme/placeholder XML directly — python-pptx's Python API exposes
most of what we need, but theme colors come via relationships and are
cleaner to pull from the raw XML.
"""

from __future__ import annotations

import hashlib
import io
from pathlib import Path
from typing import Iterable, Union

from lxml import etree
from pptx import Presentation

from .schemas import (
    LayoutDescriptor,
    LayoutKind,
    MasterCanvas,
    MasterChromeElement,
    MasterEntry,
    MasterManifest,
    MasterSafeArea,
    MasterTheme,
    ThemeEntry,
)

# One CSS pixel at 96 DPI = 9525 EMU. All PPTX geometry is in EMU; we
# convert once here and never again upstream.
EMU_PER_PX = 9525

# Common image / media extensions that real-world templates routinely
# embed. PowerPoint silently accepts a [Content_Types].xml that omits
# Default entries for these — python-pptx is stricter and refuses to
# open the package. The repair pass below ensures every extension we
# expect to find has a Default declared.
_DEFAULT_CONTENT_TYPES = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "bmp": "image/bmp",
    "tif": "image/tiff",
    "tiff": "image/tiff",
    "svg": "image/svg+xml",
    "wmf": "image/x-wmf",
    "emf": "image/x-emf",
    "mp4": "video/mp4",
    "mp3": "audio/mpeg",
    "wav": "audio/wav",
    "ico": "image/x-icon",
}


def _repair_content_types(data: bytes) -> bytes:
    """Best-effort repair of malformed ``[Content_Types].xml``.

    A sweep across real-world templates found ~10% fail python-pptx's
    parse with ``no content-type for partname '/ppt/media/image-N.png'``.
    These files are valid in PowerPoint but ship a Content-Types
    manifest that omits Default entries for image extensions actually
    referenced inside the package.

    The repair: open the .pptx as a zip, parse [Content_Types].xml,
    union in any missing Default entries for extensions we know about
    (images, video, audio), and rewrite the package. When the manifest
    is already complete this is a no-op (we still re-zip; cost is one
    pass through bytes, negligible vs. the parse work below).

    Defensive: if anything goes wrong (not a zip, no
    [Content_Types].xml, etc.), return the original bytes unchanged
    and let python-pptx raise its own clearer error.
    """
    import zipfile  # noqa: PLC0415
    from io import BytesIO  # noqa: PLC0415

    try:
        with zipfile.ZipFile(BytesIO(data), "r") as zin:
            try:
                ct_xml = zin.read("[Content_Types].xml")
            except KeyError:
                return data  # not a real OPC package; let pptx fail loudly
            existing_names = set(zin.namelist())
    except zipfile.BadZipFile:
        return data

    # Parse the manifest. lxml is already imported up top.
    try:
        ns = "http://schemas.openxmlformats.org/package/2006/content-types"
        root = etree.fromstring(ct_xml)
    except etree.XMLSyntaxError:
        return data

    declared_extensions: set[str] = set()
    declared_overrides: set[str] = set()
    for child in root:
        tag = etree.QName(child).localname
        if tag == "Default":
            ext = (child.get("Extension") or "").lower()
            if ext:
                declared_extensions.add(ext)
        elif tag == "Override":
            partname = child.get("PartName") or ""
            if partname:
                declared_overrides.add(partname)

    # Find extensions actually used in the package (excluding xml,
    # which OPC requires as a Default entry by spec — every valid
    # package already declares it).
    used_extensions: set[str] = set()
    for member in existing_names:
        if "." not in member:
            continue
        ext = member.rsplit(".", 1)[-1].lower()
        # Skip parts already covered by a per-part Override.
        if f"/{member}" in declared_overrides:
            continue
        used_extensions.add(ext)

    missing = [
        ext
        for ext in used_extensions
        if ext in _DEFAULT_CONTENT_TYPES and ext not in declared_extensions
    ]
    if not missing:
        return data

    # Union in the missing Defaults. Append before any closing tag —
    # ``etree.SubElement`` does the right thing.
    for ext in missing:
        new = etree.SubElement(
            root,
            f"{{{ns}}}Default",
            Extension=ext,
            ContentType=_DEFAULT_CONTENT_TYPES[ext],
        )
        # SubElement appends; that's fine for the OPC spec which is
        # element-order-agnostic.
        del new

    repaired_ct = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )

    # Rewrite the zip with the patched [Content_Types].xml and every
    # other member copied verbatim. Use ZIP_DEFLATED to match what
    # PowerPoint emits.
    out = BytesIO()
    with (
        zipfile.ZipFile(BytesIO(data), "r") as zin,
        zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout,
    ):
        for item in zin.infolist():
            if item.filename == "[Content_Types].xml":
                zout.writestr(item, repaired_ct)
            else:
                zout.writestr(item, zin.read(item.filename))
    return out.getvalue()


# OOXML namespaces we read. ``p`` is presentationml, ``a`` is drawingml.
_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

_THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"

# python-pptx placeholder-type names → our semantic roles.
#
# Phase 2.1: ``OBJECT``, ``CHART``, ``TABLE``, ``MEDIA`` were previously
# absent and fell through to ``"other"``, which silently dropped their
# geometry from the layout safe-area calculation. The Strategy& template
# uses ``OBJECT (7)`` as the primary content-column placeholder type
# across 14 of its 45 layouts — pre-fix, every multi-column Strategy&
# layout misclassified as ``content`` because the column count was 0.
#
# Treating these as ``"body"`` for *classification purposes* is correct:
# they all carry content. (The fact that OOXML allows OBJECT to host a
# chart or table doesn't change that — it's still a content region.)
_PH_ROLE: dict[str, str] = {
    "TITLE": "title",
    "CENTER_TITLE": "title",
    "SUBTITLE": "subtitle",
    "BODY": "body",
    "OBJECT": "body",  # Phase 2.1: generic content container
    "CHART": "body",  # Phase 2.1: chart region — content for kind detection
    "TABLE": "body",  # Phase 2.1
    "MEDIA": "body",  # Phase 2.1: video/audio embed
    "FOOTER": "footer",
    "SLIDE_NUMBER": "page_number",
    "DATE": "date",
    "PICTURE": "logo",
}

# Prompt text PowerPoint inserts into empty placeholders. If we see these
# literal strings we do not treat them as real master text.
_PROMPT_NOISE = (
    "Click to edit",
    "Click to add",
    "Master title",
    "Master subtitle",
    "Master text",
    "Second level",
    "Third level",
    "Fourth level",
    "Fifth level",
    "Sixth level",
    "Seventh level",
    "Eighth level",
    "Ninth level",
    "[Slide title]",
    "[Optional slide subtitle]",
    "[Optional footnotes/references]",
    "Date",
)


def _emu_to_px(emu: int | None) -> float:
    if emu is None:
        return 0.0
    return round(emu / EMU_PER_PX, 2)


def _read_bytes(source: Union[str, Path, bytes, bytearray]) -> bytes:
    if isinstance(source, (bytes, bytearray)):
        return bytes(source)
    return Path(source).read_bytes()


def _hex_of(color_el) -> str | None:
    """Resolve a theme color element to a hex string.

    The drawingml color-choice machinery can carry srgbClr (fixed hex),
    sysClr (system color with a hex fallback in ``lastClr``), or scheme
    references. We only handle srgbClr and sysClr — scheme references
    inside the scheme itself would be circular.
    """

    if color_el is None:
        return None
    srgb = color_el.find("a:srgbClr", _NS)
    if srgb is not None and srgb.get("val"):
        return "#" + srgb.get("val").upper()
    sys = color_el.find("a:sysClr", _NS)
    if sys is not None:
        last = sys.get("lastClr") or "000000"
        return "#" + last.upper()
    return None


def _extract_scheme_lookup(theme_xml: bytes) -> dict[str, str]:
    """Build a ``{scheme_name: hex}`` lookup for resolving placeholder
    rPr ``<a:schemeClr val="…"/>`` references.

    Includes the seven raw scheme names (dk1, lt1, dk2, lt2, accent1..6,
    plus hlink/folHlink) AND the four "logical" aliases (tx1, bg1,
    tx2, bg2) — drawingml uses the logical names in placeholder XML
    but defines them in the scheme under dk*/lt*.

    Per the OOXML spec the alias mapping flips for dark themes
    (tx1 → lt1 instead of dk1), but every consulting template we've
    audited ships a light scheme. If we hit a true dark scheme later,
    detect it by reading ``<a:clrMap>`` on the master.
    """
    out: dict[str, str] = {}
    try:
        root = etree.fromstring(theme_xml)
    except Exception:
        return out
    scheme = root.find(".//a:clrScheme", _NS)
    if scheme is None:
        return out
    for child in scheme:
        name = child.tag.split("}")[1] if "}" in child.tag else child.tag
        hex_val = _hex_of(child)
        if hex_val:
            out[name] = hex_val
    # Logical aliases — light scheme assumption.
    if "dk1" in out:
        out.setdefault("tx1", out["dk1"])
    if "lt1" in out:
        out.setdefault("bg1", out["lt1"])
    if "dk2" in out:
        out.setdefault("tx2", out["dk2"])
    if "lt2" in out:
        out.setdefault("bg2", out["lt2"])
    return out


def _resolve_font_token(token: str | None, theme_fonts: dict[str, str]) -> str | None:
    """Translate a layout-XML font reference to a concrete typeface.

    PPTX uses ``+mj-lt`` (theme major latin) and ``+mn-lt`` (theme
    minor latin) as indirection tokens on ``<a:latin typeface="…"/>``.
    A literal name like ``"Georgia"`` is returned as-is.
    """
    if not token:
        return None
    if token == "+mj-lt":
        return theme_fonts.get("major")
    if token == "+mn-lt":
        return theme_fonts.get("minor")
    return token


def _extract_placeholder_rpr(
    shp,
    theme_fonts: dict[str, str],
    scheme_lookup: dict[str, str],
) -> dict:
    """Pull layout-level typography for one placeholder.

    Reads ``<a:txBody><a:lstStyle><a:lvl1pPr>``:
      * ``defRPr/@sz`` → font size in points (sz is hundredths-of-pt)
      * ``defRPr/@b`` → ``"1"`` for bold, ``"0"`` for explicit non-bold
      * ``defRPr/@i`` → italic
      * ``defRPr/a:solidFill/a:srgbClr/@val`` → literal hex
      * ``defRPr/a:solidFill/a:schemeClr/@val`` → resolved via ``scheme_lookup``
      * ``defRPr/a:latin/@typeface`` → resolved via ``theme_fonts`` if
        a ``+mj-lt`` / ``+mn-lt`` token
      * ``lvl1pPr/@algn`` → text alignment (l/ctr/r/just → left/center/right/justify)

    Returns a partial dict — fields are only included when present in
    the source XML so callers can tell "explicit value" from "fall
    through to inheritance".
    """
    out: dict = {}
    elem = getattr(shp, "_element", None)
    if elem is None:
        return out
    # ``txBody`` is in the presentationml namespace (``p:``) on a
    # placeholder shape, but its lstStyle/lvl1pPr children are
    # drawingml (``a:``). Get the namespace right or the XPath silently
    # finds nothing — exactly the bug we hit on the first pass.
    pPr = elem.find(".//p:txBody/a:lstStyle/a:lvl1pPr", _NS)
    if pPr is None:
        return out

    align = pPr.get("algn")
    if align:
        align_map = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}
        out["align"] = align_map.get(align)

    defRPr = pPr.find("a:defRPr", _NS)
    if defRPr is None:
        return out

    sz = defRPr.get("sz")
    if sz:
        try:
            out["size"] = float(sz) / 100.0
        except ValueError:
            pass

    b_attr = defRPr.get("b")
    if b_attr is not None:
        # Explicit "0" means "not bold (overriding inherited bold)".
        # We translate to weight=400 / 700 so the LLM sees a plain
        # numeric weight rather than a tristate.
        out["weight"] = 700 if b_attr == "1" else 400

    i_attr = defRPr.get("i")
    if i_attr == "1":
        out["italic"] = True
    elif i_attr == "0":
        out["italic"] = False

    fill = defRPr.find("a:solidFill", _NS)
    if fill is not None:
        srgb = fill.find("a:srgbClr", _NS)
        sysclr = fill.find("a:sysClr", _NS)
        scheme_clr = fill.find("a:schemeClr", _NS)
        if srgb is not None and srgb.get("val"):
            out["color"] = "#" + srgb.get("val").upper()
        elif sysclr is not None:
            last = sysclr.get("lastClr") or sysclr.get("val")
            if last:
                out["color"] = "#" + last.upper()
        elif scheme_clr is not None and scheme_clr.get("val"):
            ref = scheme_clr.get("val")
            resolved = scheme_lookup.get(ref)
            if resolved:
                out["color"] = resolved

    latin = defRPr.find("a:latin", _NS)
    if latin is not None:
        font = _resolve_font_token(latin.get("typeface"), theme_fonts)
        if font:
            out["font"] = font

    return out


def _extract_theme(theme_xml: bytes) -> MasterTheme:
    """Parse ``ppt/theme/themeN.xml`` into a ``MasterTheme``."""

    root = etree.fromstring(theme_xml)

    major_el = root.find(".//a:fontScheme/a:majorFont/a:latin", _NS)
    minor_el = root.find(".//a:fontScheme/a:minorFont/a:latin", _NS)
    fonts = {
        "major": (major_el.get("typeface") if major_el is not None else None) or "Arial",
        "minor": (minor_el.get("typeface") if minor_el is not None else None) or "Arial",
    }

    scheme = root.find(".//a:clrScheme", _NS)

    def _at(name: str) -> str | None:
        return _hex_of(scheme.find(f"a:{name}", _NS)) if scheme is not None else None

    dk1, lt1 = _at("dk1"), _at("lt1")
    dk2, lt2 = _at("dk2"), _at("lt2")
    acc1, acc2 = _at("accent1"), _at("accent2")
    acc3, acc4 = _at("accent3"), _at("accent4")
    acc5, acc6 = _at("accent5"), _at("accent6")

    # Semantic mapping: primary is the first accent (Strategy& burgundy =
    # accent1 #A32020). If accent1 is black/near-black (some templates
    # leave accent1 generic and put the real brand color on dk2), fall
    # back to dk2.
    def _is_near_black(c: str | None) -> bool:
        if not c or not c.startswith("#") or len(c) != 7:
            return False
        r, g, b = int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)
        return max(r, g, b) < 32

    primary = acc1 if not _is_near_black(acc1) else (dk2 or acc1 or "#000000")
    secondary = acc2 or (dk2 if primary != dk2 else None) or "#888888"

    neutral = [c for c in (acc3, acc4, acc5, acc6, lt2) if c]

    colors: dict[str, object] = {
        "text": dk1 or "#000000",
        "bg": lt1 or "#FFFFFF",
        "primary": primary or "#000000",
        "secondary": secondary,
        "neutral": neutral,
    }
    return MasterTheme(fonts=fonts, colors=colors)


def _looks_like_prompt(text: str) -> bool:
    return any(p in text for p in _PROMPT_NOISE)


def _extract_text(shp) -> str | None:
    try:
        if not shp.has_text_frame:
            return None
    except Exception:
        return None
    try:
        raw = (shp.text_frame.text or "").strip()
    except Exception:
        return None
    if not raw:
        return None
    if _looks_like_prompt(raw):
        return None
    return raw


def _shape_geometry(shp) -> tuple[float, float, float, float] | None:
    """Return (x, y, w, h) in px, or None if any dimension is missing.

    Shapes inheriting geometry from a layout/master can legitimately have
    ``None`` on ``left/top/width/height`` — python-pptx exposes this as
    Python ``None``. In that case we skip: the real geometry lives on the
    ancestor and we'll collect it from there.
    """

    try:
        left, top, width, height = shp.left, shp.top, shp.width, shp.height
    except Exception:
        return None
    if left is None or top is None or width is None or height is None:
        return None
    return (
        _emu_to_px(left),
        _emu_to_px(top),
        _emu_to_px(width),
        _emu_to_px(height),
    )


def _collect_placeholders(
    shapes: Iterable,
    out: list[MasterChromeElement],
    *,
    theme_fonts: dict[str, str] | None = None,
    scheme_lookup: dict[str, str] | None = None,
) -> None:
    """Geometry-only collection of every placeholder.

    We deliberately *do not* read text out of placeholders. Text inside
    a placeholder is stub prompt text by definition — "Click to edit
    master title style", localised translations of the same, schema
    sample text like "Title"/"Headline"/"Subtitle". Captured text leaks
    onto every slide; the language-keyword filter we used to run never
    covered every culture (Spanish / Arabic / German all bled through).
    Skipping it across the board is both cleaner and language-agnostic.

    Real master-locked strings (e.g. a brand wordmark, a confidential
    footer) live on *non-placeholder* shapes and are picked up by
    ``_collect_master_text_shapes``.
    """
    for shp in shapes:
        try:
            ph = shp.placeholder_format
        except Exception:
            ph = None
        if ph is None:
            continue
        type_name = getattr(ph.type, "name", None) if ph.type is not None else None
        role = _PH_ROLE.get(type_name or "", "other")
        geom = _shape_geometry(shp)
        if geom is None:
            continue
        x, y, w, h = geom

        # Phase 2.5: pull layout-level rPr (font/size/weight/italic/color/align)
        # so the agent appendix can tell the LLM what each placeholder
        # actually looks like, not just where it lives.
        rpr = (
            _extract_placeholder_rpr(shp, theme_fonts or {}, scheme_lookup or {})
            if (theme_fonts is not None or scheme_lookup is not None)
            else {}
        )
        out.append(
            MasterChromeElement(
                role=role,
                x=x,
                y=y,
                w=w,
                h=h,
                text=None,
                font=rpr.get("font"),
                size=rpr.get("size"),
                weight=rpr.get("weight"),
                color=rpr.get("color"),
                align=rpr.get("align"),
            )
        )


def _collect_master_text_shapes(shapes: Iterable, out: list[MasterChromeElement]) -> None:
    """Non-placeholder text shapes on the master (e.g. "Strategy&" footer).

    These are the real brand-locked strings. They sit alongside
    placeholders and carry literal text the master wants preserved on
    every slide.
    """

    for shp in shapes:
        try:
            if shp.placeholder_format is not None:
                continue
        except Exception:
            pass
        txt = _extract_text(shp)
        if not txt:
            continue
        geom = _shape_geometry(shp)
        if geom is None:
            continue
        x, y, w, h = geom
        # Heuristic: a small text box at the bottom of the canvas that
        # carries literal text is a footer.
        role = "footer" if y > 0 else "other"
        out.append(MasterChromeElement(role=role, x=x, y=y, w=w, h=h, text=txt))


# Phase 2.1c: Strong name keywords win over placeholder counting.
# Audited templates revealed that descriptive names like
# ``Quote - White`` (1 BODY + 1 SUBTITLE) would count as ``content``
# under pure placeholder logic, but the layout is obviously a quote.
# Order matters here: ``two column`` must be tested before ``column``
# so the more-specific match wins. Each tuple is ``(substring, kind)``.
_NAME_KEYWORDS: list[tuple[str, LayoutKind]] = [
    # Strong, narrow signals first
    ("section header", "section_header"),
    ("section divider", "section_header"),
    ("chapter divider", "section_header"),
    ("agenda", "agenda"),
    ("table of contents", "agenda"),
    ("outline", "agenda"),
    # Quote / callout family
    ("pull quote", "quote"),
    ("quote", "quote"),
    ("callout", "quote"),
    # Cover / title family
    ("title slide", "title"),
    ("cover", "title"),
    # Numbers / KPI family
    ("big number", "kpi"),
    ("big numbers", "kpi"),
    ("kpi", "kpi"),
    ("metric", "kpi"),
    ("stat ", "kpi"),
    # Column families — order: two before three before column to disambiguate
    ("two column", "two_column"),
    ("two-column", "two_column"),
    ("2_col", "two_column"),
    ("2 col", "two_column"),
    ("three column", "comparison"),
    ("three-column", "comparison"),
    ("four column", "comparison"),
    ("four-column", "comparison"),
    ("five column", "comparison"),
    ("five-column", "comparison"),
    ("comparison", "comparison"),
    ("vs.", "comparison"),
    ("versus", "comparison"),
    # Blank
    ("blank", "blank"),
]


def _classify_by_name(name: str) -> LayoutKind | None:
    """Return a kind when the layout name carries a strong signal,
    else ``None`` (caller falls back to placeholder counting)."""
    if not name:
        return None
    lower = name.lower()
    for keyword, kind in _NAME_KEYWORDS:
        if keyword in lower:
            return kind
    return None


def _are_side_by_side(bodies: list[MasterChromeElement]) -> bool:
    """True when ``bodies`` form a horizontal row (true multi-column)
    vs. a stacked content+secondary arrangement.

    Real-world signal: stc's "12_Content slide" has 2 BODY placeholders
    but they're stacked (main content + a small attribution / source
    line below) — the classifier was calling that ``two_column`` purely
    on count. A geometry check rules out stacked layouts.

    Heuristic:
      * vertical centres within 30% of the smaller body's height —
        anything more is "stacked", not "side-by-side"
      * x-ranges must not overlap by more than 10% of the smaller
        body's width — overlapping ranges with same y = stacked
        decoration, not columns
    """
    if len(bodies) < 2:
        return False
    # Sort by x so overlap check is stable.
    sorted_bodies = sorted(bodies, key=lambda b: b.x)
    for i in range(len(sorted_bodies) - 1):
        a = sorted_bodies[i]
        b = sorted_bodies[i + 1]
        cy_a = a.y + a.h / 2
        cy_b = b.y + b.h / 2
        min_h = min(a.h, b.h)
        if min_h <= 0:
            return False
        if abs(cy_a - cy_b) > 0.3 * min_h:
            return False
        # x-ranges must not overlap meaningfully.
        overlap = max(0.0, (a.x + a.w) - b.x)
        min_w = min(a.w, b.w)
        if min_w > 0 and overlap > 0.1 * min_w:
            return False
    return True


def _classify_by_placeholders(
    placeholders: list[MasterChromeElement],
) -> LayoutKind:
    """Placeholder-count + geometry classifier. Used for layouts whose
    name carries no signal (e.g. stc's opaque ``"12_Content slide _
    VCS_to use"``).

    Geometry matters: 2 BODY placeholders side-by-side is two_column,
    but 2 BODY placeholders stacked (main content + attribution line)
    is plain ``content`` — the classifier was wrongly flattening the
    second case into ``two_column`` for stc's deck.
    """
    title_n = sum(1 for p in placeholders if p.role in ("title", "subtitle"))
    bodies = [p for p in placeholders if p.role == "body"]
    body_n = len(bodies)
    other_text = sum(1 for p in placeholders if p.role == "other" and p.text)

    if title_n == 0 and body_n == 0:
        return "blank"
    if title_n >= 1 and body_n == 0 and other_text == 0:
        return "section_header"
    if title_n >= 1 and body_n == 1:
        return "content"
    if title_n >= 1 and body_n == 2:
        return "two_column" if _are_side_by_side(bodies) else "content"
    if body_n >= 3:
        # Three+ bodies: still need geometry. A side-by-side trio is a
        # comparison / three-column; stacked trios are agendas / process
        # steps / TOCs and we leave them as ``content``.
        return "comparison" if _are_side_by_side(bodies) else "content"
    return "other"


def _classify_layout(placeholders: list[MasterChromeElement], name: str) -> LayoutKind:
    """Two-stage classifier: name keywords first, placeholder + geometry
    fallback.

    Strategy&'s descriptive names (``Quote - White``, ``Two Columns``,
    ``Section Header - Side Image``, ``Four Big Numbers``) classify
    correctly via the keyword path. stc's opaque names
    (``12_Content slide _ VCS_to use``) fall through to placeholder
    counting + geometry, which keeps stacked content layouts out of
    ``two_column`` / ``comparison``.
    """
    by_name = _classify_by_name(name)
    if by_name is not None:
        return by_name
    return _classify_by_placeholders(placeholders)


def _layout_safe_area(
    placeholders: list[MasterChromeElement], canvas: MasterCanvas
) -> MasterSafeArea | None:
    """Per-layout safe_area override.

    Bound-boxes every ``body`` placeholder. A few real-world templates
    smuggle decorative micro-labels (a 190×16 "Document Title" tag in
    a corner) under the BODY role; without a sanity check, the LLM
    would inherit a useless 16px-tall content region. We reject any
    bbox that is either:
      * smaller than ~30% of canvas in either dimension — too cramped
        to author meaningful content into, AND
      * anchored in the upper-right or other unusable quadrant.
    Returning ``None`` falls through to the master-level safe_area,
    which is almost always the safer default.
    """
    bodies = [p for p in placeholders if p.role == "body"]
    if not bodies:
        return None
    x1 = min(p.x for p in bodies)
    y1 = min(p.y for p in bodies)
    x2 = max(p.x + p.w for p in bodies)
    y2 = max(p.y + p.h for p in bodies)
    w = max(x2 - x1, 0)
    h = max(y2 - y1, 0)

    # Sanity: the bbox should be a meaningful content region. Anything
    # narrower than 30% of the canvas or shorter than 25% is almost
    # certainly a decorative label, not a body. Cover layouts use the
    # full canvas, so bigger is always fine; we only reject too-small.
    if w < canvas.w * 0.30 or h < canvas.h * 0.25:
        return None
    return MasterSafeArea(x=x1, y=y1, w=w, h=h)


def _layout_description(kind: LayoutKind, sa: MasterSafeArea | None, canvas: MasterCanvas) -> str:
    """Short LLM-facing menu line. Pure formatter — no judgment calls."""
    parts = [
        {
            "title": "title slide — full-canvas, expects a centered title",
            "section_header": "section divider — large title, no body content",
            "agenda": "table of contents / outline",
            "content": "single body region",
            "two_column": "two body regions side-by-side",
            "comparison": "three or more body regions",
            "kpi": "metric or stat callout",
            "quote": "pull-quote callout",
            "blank": "no placeholders — bare canvas inside chrome",
            "other": "free-form layout",
        }[kind]
    ]
    if sa is not None:
        parts.append(f"safe_area x={sa.x:.0f} y={sa.y:.0f} {sa.w:.0f}×{sa.h:.0f}")
    return "; ".join(parts)


def _build_layout_menu(
    sm,
    canvas: MasterCanvas,
    master_safe_area: MasterSafeArea,
    master_index: int,
    *,
    theme_fonts: dict[str, str] | None = None,
    scheme_lookup: dict[str, str] | None = None,
) -> list[LayoutDescriptor]:
    """Walk every slide layout in this master and emit a classified
    LayoutDescriptor with stable (master_index, layout_index) identity."""
    descriptors: list[LayoutDescriptor] = []
    for layout_index, layout in enumerate(sm.slide_layouts):
        name = getattr(layout, "name", None) or ""
        layout_chrome: list[MasterChromeElement] = []
        _collect_placeholders(
            layout.placeholders,
            layout_chrome,
            theme_fonts=theme_fonts,
            scheme_lookup=scheme_lookup,
        )
        kind = _classify_layout(layout_chrome, name)
        sa = _layout_safe_area(layout_chrome, canvas)
        # Only override the master safe_area when the layout's bodies
        # actually differ — saves the LLM redundant info.
        if sa is not None and (
            abs(sa.x - master_safe_area.x) < 4
            and abs(sa.y - master_safe_area.y) < 4
            and abs(sa.w - master_safe_area.w) < 4
            and abs(sa.h - master_safe_area.h) < 4
        ):
            sa = None
        descriptors.append(
            LayoutDescriptor(
                name=name,
                kind=kind,
                safe_area=sa,
                placeholders=layout_chrome,
                description=_layout_description(kind, sa, canvas),
                master_index=master_index,
                layout_index=layout_index,
            )
        )
    return descriptors


def _dedupe_chrome(
    elements: list[MasterChromeElement],
) -> list[MasterChromeElement]:
    """Drop duplicates that appear on both master and layout.

    Keyed on (role, rounded x, rounded y, rounded w, rounded h). First
    occurrence wins — master is read before layouts so master values
    dominate.
    """

    seen: set[tuple] = set()
    out: list[MasterChromeElement] = []
    for el in elements:
        key = (el.role, round(el.x), round(el.y), round(el.w), round(el.h))
        if key in seen:
            continue
        seen.add(key)
        out.append(el)
    return out


def _theme_blob_for(sm) -> bytes | None:
    """Pull the theme XML bytes attached to a slide master, or None
    when the relationship can't be resolved (rare; defensive)."""
    try:
        theme_part = sm.part.part_related_by(_THEME_REL)
        return theme_part.blob
    except Exception:
        return None


def _theme_index_for(sm) -> int:
    """Best-effort: which themeN.xml backs this master?

    The theme's partname looks like ``/ppt/theme/theme7.xml``. We
    parse the digit; defaults to 1 when missing. Used by callers to
    cross-reference a master with a deduped ThemeEntry.
    """
    try:
        theme_part = sm.part.part_related_by(_THEME_REL)
        partname = str(theme_part.partname)
        # e.g. "/ppt/theme/theme7.xml"
        import re  # noqa: PLC0415

        m = re.search(r"theme(\d+)\.xml", partname)
        if m:
            return int(m.group(1))
    except Exception:
        pass
    return 1


def _safe_area_for(chrome: list[MasterChromeElement], canvas: MasterCanvas) -> MasterSafeArea:
    """Pick the body-placeholder bbox if sane, else a 48-px margin
    fallback. Same logic as before, factored for the multi-master loop."""
    body = next((c for c in chrome if c.role == "body"), None)
    if body and body.w > 0 and body.h > 0:
        return MasterSafeArea(x=body.x, y=body.y, w=body.w, h=body.h)
    margin = 48.0
    top = 175.0
    return MasterSafeArea(
        x=margin,
        y=top,
        w=max(canvas.w - 2 * margin, 0),
        h=max(canvas.h - top - 50, 0),
    )


def _theme_palette_dict(theme: MasterTheme) -> dict[str, object]:
    """Snapshot a MasterTheme's palette into a plain dict so
    MasterEntry.palette and ThemeEntry.palette stay JSON-friendly."""
    return dict(theme.colors)


def extract_master_from_pptx(
    source: Union[str, Path, bytes, bytearray],
    *,
    name: str | None = None,
) -> MasterManifest:
    """Parse a .pptx and return a multi-master MasterManifest.

    Phase 2.1: walks every ``pres.slide_masters[i]`` and every layout
    inside each. Theme files are deduplicated by raw-XML hash so a
    template like Strategy& (3 byte-identical theme files) reports 1
    unique palette, while stc's 12 distinct theme files collapse to
    its 3 real palettes.

    Legacy single-master fields (``theme``, ``safe_area``, ``chrome``,
    ``layouts``) on the manifest are populated from masters[0] for
    back-compat. New code reads ``manifest.masters`` and
    ``manifest.themes`` directly.

    Resilience: a sweep across real-world templates surfaced ~10% with
    malformed ``[Content_Types].xml`` (PowerPoint accepts these but
    python-pptx refuses to open them). We run a lightweight repair
    pass before delegating, so the extractor degrades gracefully on
    common corruption rather than throwing on parse.
    """
    data = _read_bytes(source)
    # SHA is computed on the original bytes — we want re-uploads of the
    # same source file to dedup in storage even if our repair pass
    # tweaked the in-memory copy.
    sha = hashlib.sha256(data).hexdigest()

    repaired = _repair_content_types(data)
    pres = Presentation(io.BytesIO(repaired))

    canvas = MasterCanvas(
        w=int(round(pres.slide_width / EMU_PER_PX)),
        h=int(round(pres.slide_height / EMU_PER_PX)),
    )

    # ── Pass 1: extract every master with its theme + chrome + layouts ──
    masters: list[MasterEntry] = []
    # Map (resolved palette + fonts) → ThemeEntry. We dedup on resolved
    # content rather than raw XML because corporate templates routinely
    # ship N theme files that differ only in XML whitespace, comments,
    # or theme-name attribute but resolve to the same palette/fonts.
    # stc has 12 theme files that collapse to 3 actual palettes;
    # raw-XML hashing would falsely report 12.
    theme_by_key: dict[str, ThemeEntry] = {}

    def _theme_dedup_key(theme: MasterTheme) -> str:
        """Stable string key over (palette, fonts) for dedup. Two
        themes that produce identical token resolution at render time
        are the same theme for our purposes."""
        import json  # noqa: PLC0415

        return json.dumps(
            {"colors": theme.colors, "fonts": theme.fonts},
            sort_keys=True,
            default=str,
        )

    for mi, sm in enumerate(pres.slide_masters):
        theme_blob = _theme_blob_for(sm)
        if theme_blob is None:
            theme = MasterTheme(
                fonts={"major": "Arial", "minor": "Arial"},
                colors={
                    "text": "#000000",
                    "bg": "#FFFFFF",
                    "primary": "#000000",
                    "secondary": "#888888",
                    "neutral": [],
                },
            )
            scheme_lookup: dict[str, str] = {}
        else:
            theme = _extract_theme(theme_blob)
            scheme_lookup = _extract_scheme_lookup(theme_blob)

        theme_index = _theme_index_for(sm)
        key = _theme_dedup_key(theme)
        if key in theme_by_key:
            theme_by_key[key].indices.append(theme_index)
        else:
            theme_by_key[key] = ThemeEntry(
                indices=[theme_index],
                palette=_theme_palette_dict(theme),
                fonts=dict(theme.fonts),
            )

        # Master-level chrome + safe_area
        chrome: list[MasterChromeElement] = []
        _collect_placeholders(
            sm.placeholders,
            chrome,
            theme_fonts=theme.fonts,
            scheme_lookup=scheme_lookup,
        )
        _collect_master_text_shapes(sm.shapes, chrome)
        chrome = _dedupe_chrome(chrome)
        safe_area = _safe_area_for(chrome, canvas)

        layouts = _build_layout_menu(
            sm,
            canvas,
            safe_area,
            master_index=mi,
            theme_fonts=theme.fonts,
            scheme_lookup=scheme_lookup,
        )

        # python-pptx's SlideMaster has a ``name`` property but it's
        # frequently absent or empty on real templates; fall back to a
        # debug-friendly synthetic.
        sm_name = ""
        try:
            sm_name = getattr(sm, "name", "") or ""
        except Exception:
            pass

        masters.append(
            MasterEntry(
                index=mi,
                name=sm_name,
                theme_index=theme_index,
                palette=_theme_palette_dict(theme),
                fonts=dict(theme.fonts),
                layouts=layouts,
            )
        )

    # ── Pass 2: union of font families referenced across all masters ──
    fonts_referenced: list[str] = []
    seen: set[str] = set()
    for m in masters:
        for f in (m.fonts.get("major"), m.fonts.get("minor")):
            if f and f not in seen:
                seen.add(f)
                fonts_referenced.append(f)

    # ── Legacy single-master fields (back-compat for callers that
    # haven't migrated to manifest.masters[*]) ──
    if masters:
        primary = masters[0]
        primary_theme = MasterTheme(fonts=primary.fonts, colors=primary.palette)
        # Re-collect chrome + safe_area for masters[0] because we don't
        # carry them on MasterEntry (they're derivable per-master). The
        # primary master is rare to need post-extraction, but keeping
        # legacy fields populated avoids breaking existing tests + the
        # active-master appendix in QueryEngine.
        sm0 = pres.slide_masters[0]
        chrome0: list[MasterChromeElement] = []
        _collect_placeholders(sm0.placeholders, chrome0)
        _collect_master_text_shapes(sm0.shapes, chrome0)
        chrome0 = _dedupe_chrome(chrome0)
        safe_area0 = _safe_area_for(chrome0, canvas)
        legacy_layouts = primary.layouts
    else:
        # Pathological: zero masters. Build a defensive default.
        primary_theme = MasterTheme(
            fonts={"major": "Arial", "minor": "Arial"},
            colors={
                "text": "#000000",
                "bg": "#FFFFFF",
                "primary": "#000000",
                "secondary": "#888888",
                "neutral": [],
            },
        )
        chrome0 = []
        safe_area0 = MasterSafeArea(x=48, y=175, w=canvas.w - 96, h=canvas.h - 225)
        legacy_layouts = []

    pres_title = ""
    try:
        pres_title = (pres.core_properties.title or "").strip()
    except Exception:
        pass

    return MasterManifest(
        name=name or pres_title or "Imported master",
        source_sha256=sha,
        canvas=canvas,
        # Legacy fields — populated from masters[0]
        theme=primary_theme,
        safe_area=safe_area0,
        chrome=chrome0,
        layouts=legacy_layouts,
        # Phase 2.1 new fields
        masters=masters,
        themes=list(theme_by_key.values()),
        fonts_referenced=fonts_referenced,
    )

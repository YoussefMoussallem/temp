"""LibreOffice rendering bridge.

Two callable surfaces:

* ``build_layout_preview_pptx(pptx_bytes, [(m_idx, l_idx), ...])`` —
  pure-Python helper. Loads the source .pptx with python-pptx, appends
  one EMPTY slide per requested ``(master, layout)`` pair, returns
  the new .pptx bytes. The new slides carry no content; rendering them
  yields the LAYOUT itself (master background + layout placeholders
  + chrome) — which is what we want to show the user as a thumbnail.

* ``PptxRenderer.render_layouts`` — async client to the sidecar.
  Builds the temp .pptx, POSTs it to the renderer's ``/render``
  endpoint, parses the returned zip of PNGs into a dict keyed by
  ``(master_index, layout_index)``.

Why pre-build a temp .pptx instead of asking the sidecar to render
specific layouts: LibreOffice's ``--convert-to png`` operates on
slides, not layouts. There's no ``layout-to-png`` flag. Synthesizing
slides that *use* the layout we care about is the cleanest end-around.
"""

from __future__ import annotations

import io
import zipfile
from typing import Iterable

import httpx
from pptx import Presentation


def build_layout_preview_pptx(
    source_pptx: bytes,
    specs: Iterable[tuple[int, int]],
    *,
    strip_existing_slides: bool = False,
) -> bytes:
    """Append one new (empty) slide per ``(master_index, layout_index)``
    spec and return the resulting .pptx bytes.

    By default the source .pptx's existing slides are preserved. Pass
    ``strip_existing_slides=True`` to drop them first — useful when
    the caller wants ONLY the synthesised previews in the rendered
    output (so the sidecar's slide indices map cleanly to the spec
    list, with no source slides in the way).

    Raises ``IndexError`` when a spec references a master or layout
    that doesn't exist; the upload endpoint catches and surfaces 422.
    """
    spec_list = list(specs)
    src = Presentation(io.BytesIO(source_pptx))

    # Validate up front so we don't mutate halfway and fail mid-flight.
    masters = list(src.slide_masters)
    for mi, li in spec_list:
        if mi < 0 or mi >= len(masters):
            raise IndexError(f"master_index {mi} out of range (deck has {len(masters)} masters)")
        layouts = list(masters[mi].slide_layouts)
        if li < 0 or li >= len(layouts):
            raise IndexError(
                f"layout_index {li} out of range for master {mi} "
                f"(master has {len(layouts)} layouts)"
            )

    if strip_existing_slides:
        _strip_slides(src)

    # Append the new slides.
    for mi, li in spec_list:
        layout = list(masters[mi].slide_layouts)[li]
        src.slides.add_slide(layout)

    buf = io.BytesIO()
    src.save(buf)
    return buf.getvalue()


def _strip_slides(pres) -> None:
    """Detach every slide from the Presentation so add_slide() can
    re-issue ``slide1.xml`` cleanly.

    Removes the sldId entries from the presentation XML *and* drops
    the slide relationships from the presentation part. Orphan slide
    parts then prune themselves on ``pres.save()`` because python-pptx
    walks reachable parts during package serialisation. Without
    dropping the rels, ``next_partname_idx`` would still see slide1
    .xml in the package and collide.
    """
    sldIdLst = pres.slides._sldIdLst
    rels = pres.part.rels
    rId_attrib = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for sldId in list(sldIdLst):
        rId = sldId.attrib.get(rId_attrib)
        sldIdLst.remove(sldId)
        if rId and rId in rels:
            try:
                pres.part.drop_rel(rId)
            except KeyError:
                # Already dropped or a stale id; either way, nothing left to do.
                pass


def _count_slides(pptx_bytes: bytes) -> int:
    """How many slides are already in ``pptx_bytes``."""
    return len(list(Presentation(io.BytesIO(pptx_bytes)).slides))


def find_layout_representative_slides(
    pptx_bytes: bytes,
) -> dict[tuple[int, int], int]:
    """Map each (master_index, layout_index) to a 1-based slide
    index in the deck that *uses* that layout.

    Real corporate templates routinely ship demo decks where the
    author has already placed an example slide on every interesting
    layout — those examples are gold-standard previews because they
    show the layout *used* (real chrome, real photo placeholders,
    real proportions) instead of empty placeholder prompts.

    Layouts that no slide uses don't appear in the result. The
    caller is responsible for the fallback (typically: skip the
    preview and surface 'not used in source' in the curation UI).

    We match layouts by their OPC ``partname`` (e.g. ``/ppt/slideLayouts/
    slideLayout7.xml``) rather than Python identity, because python-pptx
    returns different SlideLayout proxies for the same XML part
    depending on the access path (``pres.slide_masters[mi].slide_layouts[li]``
    vs ``slide.slide_layout``). The partname is stable across both.
    """
    pres = Presentation(io.BytesIO(pptx_bytes))
    masters = list(pres.slide_masters)
    partname_to_key: dict[str, tuple[int, int]] = {}
    for mi, sm in enumerate(masters):
        for li, layout in enumerate(sm.slide_layouts):
            partname_to_key[str(layout.part.partname)] = (mi, li)

    out: dict[tuple[int, int], int] = {}
    for one_based_idx, slide in enumerate(pres.slides, start=1):
        try:
            partname = str(slide.slide_layout.part.partname)
        except Exception:
            continue
        key = partname_to_key.get(partname)
        if key is None or key in out:
            continue
        out[key] = one_based_idx
    return out


class PptxRenderer:
    """Async client for the pptx-renderer sidecar service.

    The sidecar exposes ``POST /render``; we hand it a multipart
    upload of bytes + a comma-separated list of slide indices and
    receive a zip of PNGs back.

    Construction is cheap (no I/O); instantiate ad-hoc per call —
    requests are infrequent (one batch per master upload) so a shared
    client pool isn't worth the lifecycle complexity.
    """

    def __init__(
        self,
        base_url: str | None = None,
        timeout: float = 60.0,
    ) -> None:
        # Honour PPTX_RENDERER_URL when set so host-run dev (where
        # ``pptx-renderer`` isn't a resolvable hostname) can point at
        # ``http://127.0.0.1:8002``. The docker-network hostname stays
        # the default for in-compose deployments.
        import os  # noqa: PLC0415

        self._base_url = (
            base_url or os.environ.get("PPTX_RENDERER_URL") or "http://pptx-renderer:8002"
        ).rstrip("/")
        self._timeout = timeout

    async def render_layouts(
        self,
        source_pptx: bytes,
        specs: list[tuple[int, int]],
    ) -> dict[tuple[int, int], bytes]:
        """Render the requested layouts as PNGs.

        Returns ``{(master_index, layout_index): png_bytes}``. An
        empty ``specs`` list short-circuits to ``{}`` without an HTTP
        call so the upload path can call this unconditionally even
        when there are zero layouts.

        Two-pass strategy:

        1. **Representative-slide pass.** For layouts that an actual
           slide in the deck uses, render those slides — the demo
           content shows the layout in real use, which is the most
           informative thumbnail. (Real corporate templates routinely
           ship demo decks where every used layout has an example.)
        2. **Synthesized-empty pass.** For layouts no slide uses,
           build a slide-stripped copy of the source and add one
           empty slide per dormant layout. The PNG shows master
           background + chrome + placeholder prompts ("Click to edit
           master title"). This covers pure-template uploads where
           the source has zero slides, and the long-tail of unused
           layouts in mixed decks.

        Both passes share a single HTTP round-trip when only one is
        needed; they fire two separate calls otherwise — sequential
        because the sidecar processes one .pptx per request.
        """
        if not specs:
            return {}

        layout_to_slide = find_layout_representative_slides(source_pptx)
        used_specs = [s for s in specs if s in layout_to_slide]
        unused_specs = [s for s in specs if s not in layout_to_slide]

        out: dict[tuple[int, int], bytes] = {}

        if used_specs:
            out.update(
                await self._render_representatives(
                    source_pptx,
                    used_specs,
                    layout_to_slide,
                )
            )

        if unused_specs:
            out.update(await self._render_synthesised(source_pptx, unused_specs))

        return out

    async def _render_representatives(
        self,
        source_pptx: bytes,
        wanted_specs: list[tuple[int, int]],
        layout_to_slide: dict[tuple[int, int], int],
    ) -> dict[tuple[int, int], bytes]:
        """Render layouts that an existing slide uses by plucking the
        right page out of a render of the source .pptx as-is.
        """
        slide_indices_list = sorted({layout_to_slide[s] for s in wanted_specs})
        zip_bytes = await self._render_pptx(
            source_pptx,
            slide_indices_list,
        )

        # Sidecar names PNGs ``slide-1.png``..``slide-N.png`` by
        # position-after-filter. Re-key by (m_idx, l_idx).
        idx_to_zipname = {
            slide_idx: f"slide-{position + 1}.png"
            for position, slide_idx in enumerate(slide_indices_list)
        }
        out: dict[tuple[int, int], bytes] = {}
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            for spec in wanted_specs:
                slide_idx = layout_to_slide[spec]
                zip_name = idx_to_zipname.get(slide_idx)
                if not zip_name:
                    continue
                try:
                    out[spec] = zf.read(zip_name)
                except KeyError:
                    # Sidecar dropped this slide. Skip; curation UI
                    # shows a placeholder.
                    continue
        return out

    async def _render_synthesised(
        self,
        source_pptx: bytes,
        unused_specs: list[tuple[int, int]],
    ) -> dict[tuple[int, int], bytes]:
        """Render layouts no slide uses by building a slide-stripped
        copy of the source plus one empty slide per spec.

        Spec order is preserved end-to-end: synth slide N corresponds
        to ``unused_specs[N]``, so we ask the sidecar for a page-range
        and re-key by position.
        """
        synth_pptx = build_layout_preview_pptx(
            source_pptx,
            unused_specs,
            strip_existing_slides=True,
        )
        slide_indices_list = list(range(1, len(unused_specs) + 1))
        zip_bytes = await self._render_pptx(synth_pptx, slide_indices_list)

        out: dict[tuple[int, int], bytes] = {}
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            for position, spec in enumerate(unused_specs):
                zip_name = f"slide-{position + 1}.png"
                try:
                    out[spec] = zf.read(zip_name)
                except KeyError:
                    # Synth slide didn't render — rare, but don't blow up.
                    continue
        return out

    async def _render_pptx(
        self,
        pptx_bytes: bytes,
        slide_indices: list[int],
    ) -> bytes:
        """POST .pptx + slide indices to the sidecar; return zip bytes.

        Centralised so both render passes share the same HTTP/error
        handling rather than duplicating the call site.
        """
        files = {"file": ("source.pptx", pptx_bytes, _PPTX_MIME)}
        data = {"slides": ",".join(str(i) for i in slide_indices)}
        async with httpx.AsyncClient(timeout=self._timeout) as client:
            resp = await client.post(
                f"{self._base_url}/render",
                files=files,
                data=data,
            )
            resp.raise_for_status()
            return resp.content


# ── Internal helpers ──────────────────────────────────────────────────

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _expected_zip_names(n: int) -> list[str]:
    """Sidecar names PNGs ``slide-1.png`` ... ``slide-N.png``. Kept in
    lock-step with the sidecar's writer in ``server/pptx-renderer``.
    """
    return [f"slide-{i + 1}.png" for i in range(n)]

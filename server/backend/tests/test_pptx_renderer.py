"""Phase 2.2 — LibreOffice rendering client.

Two layers:

1. **Unit** — ``build_layout_preview_pptx``: pure-Python, no I/O. Takes
   .pptx bytes + a list of (master_index, layout_index) and returns
   new .pptx bytes with exactly one new slide per requested layout.
   The new slides are *empty* on those layouts (no content), which is
   what we want — the rendered PNG shows the layout's master
   background + chrome + placeholder prompts, exactly what the user
   sees in PowerPoint's "View → Slide Master" mode.

2. **Integration** — ``PptxRenderer.render_layouts`` round-trips
   through the live sidecar at ``http://localhost:8002``. Skipped
   cleanly when the sidecar isn't reachable so a fresh checkout
   doesn't fail the suite.

Mocked HTTP tests are deliberately omitted — the wrapper around
``httpx`` is thin enough that a mock would assert the same shape
the integration test asserts. We'd rather catch real LibreOffice
behaviour drift in one well-placed integration test than mock the
SDK and find out in production.
"""

from __future__ import annotations

import io
import socket

import pytest
from pptx import Presentation


# pytestmark is intentionally NOT set — only the integration tests
# below need ``@pytest.mark.asyncio``; the unit tests are sync.


def _sidecar_reachable() -> bool:
    try:
        sock = socket.create_connection(("127.0.0.1", 8002), timeout=0.5)
        sock.close()
        return True
    except OSError:
        return False


sidecar_required = pytest.mark.skipif(
    not _sidecar_reachable(),
    reason="pptx-renderer sidecar not on 127.0.0.1:8002 — run `docker compose up -d pptx-renderer`",
)


# ── Unit: temp-pptx builder ───────────────────────────────────────────


def test_build_layout_preview_pptx_one_slide_per_spec(minimal_pptx: bytes) -> None:
    from app.services.pptx_renderer import build_layout_preview_pptx

    # Default Office minimal has 1 master with 11 layouts. Ask for
    # layouts 0, 4, 7 from master 0; expect exactly 3 new slides.
    out_bytes = build_layout_preview_pptx(minimal_pptx, [(0, 0), (0, 4), (0, 7)])
    out_pres = Presentation(io.BytesIO(out_bytes))
    # The original .pptx has zero slides; we add 3, expect 3.
    assert len(out_pres.slides) == 3


def test_build_layout_preview_pptx_orders_slides_by_spec(minimal_pptx: bytes) -> None:
    """Order matters: caller indexes the returned PNGs by position
    in the original spec list, so slide N must correspond to spec[N]."""
    from app.services.pptx_renderer import build_layout_preview_pptx

    out_bytes = build_layout_preview_pptx(
        minimal_pptx,
        [(0, 1), (0, 0)],  # reversed
    )
    out_pres = Presentation(io.BytesIO(out_bytes))
    assert len(out_pres.slides) == 2
    # Slide 0 should use the layout we asked for first (layout 1)
    # The slide_layout property exposes the layout instance; comparing
    # ids would be brittle (python-pptx doesn't expose a stable id),
    # but layouts within a master are ordered, so we can compare
    # names instead.
    src = Presentation(io.BytesIO(minimal_pptx))
    expected_first = src.slide_masters[0].slide_layouts[1].name
    expected_second = src.slide_masters[0].slide_layouts[0].name
    assert out_pres.slides[0].slide_layout.name == expected_first
    assert out_pres.slides[1].slide_layout.name == expected_second


def test_build_layout_preview_pptx_rejects_invalid_indices(minimal_pptx: bytes) -> None:
    from app.services.pptx_renderer import build_layout_preview_pptx

    with pytest.raises(IndexError):
        build_layout_preview_pptx(minimal_pptx, [(99, 0)])  # no master 99
    with pytest.raises(IndexError):
        build_layout_preview_pptx(minimal_pptx, [(0, 999)])  # no layout 999


def test_build_layout_preview_pptx_strips_existing_slides(
    minimal_pptx_with_slides: bytes,
) -> None:
    """Phase 2.4.1: when a source deck has slides already, the strip
    flag drops them so the synthesised previews start at slide 1.

    Without stripping, python-pptx's ``add_slide`` collides on
    ``slide1.xml`` partnames and produces a corrupt zip. With
    stripping, the source slides are detached from the presentation
    XML and the new ones get clean partnames.
    """
    from app.services.pptx_renderer import build_layout_preview_pptx

    src = Presentation(io.BytesIO(minimal_pptx_with_slides))
    assert len(src.slides) > 0  # fixture sanity
    out_bytes = build_layout_preview_pptx(
        minimal_pptx_with_slides,
        [(0, 0), (0, 1)],
        strip_existing_slides=True,
    )
    out_pres = Presentation(io.BytesIO(out_bytes))
    # Only the two synthesised slides should remain.
    assert len(out_pres.slides) == 2
    # And the masters/layouts are intact.
    assert len(out_pres.slide_masters) == len(src.slide_masters)


def test_find_layout_representative_slides_minimal(minimal_pptx: bytes) -> None:
    """A bare-empty deck has zero slides and therefore zero
    representatives. Synthetic minimal proves the function survives
    the no-slides edge case without errors."""
    from app.services.pptx_renderer import find_layout_representative_slides

    assert find_layout_representative_slides(minimal_pptx) == {}


def test_find_layout_representative_slides_real_deck(private_pptx_path) -> None:
    """A real corporate deck (Strategy&) uses 41 distinct layouts
    across its 49 slides. The representative-slides map should
    contain one entry per used layout, indexed at the FIRST slide
    that uses it. Skipped without the private fixture."""
    from app.services.pptx_renderer import find_layout_representative_slides

    path = private_pptx_path("strategy-and-white")
    data = path.read_bytes()
    reps = find_layout_representative_slides(data)

    # Audit: 41/45 layouts are used. We expect one entry per used
    # layout, and every entry is a 1-based slide index in [1..49].
    assert 35 <= len(reps) <= 45
    for slide_idx in reps.values():
        assert 1 <= slide_idx <= 49


# ── Integration: live sidecar round-trip ──────────────────────────────


@sidecar_required
@pytest.mark.asyncio
async def test_render_layouts_round_trip_against_sidecar(
    minimal_pptx_with_slides: bytes,
) -> None:
    """Drive the full path: representative-slide lookup → render via
    sidecar → parse zip → return PNG dict. Uses the with-slides
    fixture because the renderer's strategy needs real slides on the
    requested layouts.
    """
    from app.services.pptx_renderer import PptxRenderer

    renderer = PptxRenderer(base_url="http://127.0.0.1:8002")
    specs = [(0, 0), (0, 1), (0, 2)]
    pngs = await renderer.render_layouts(minimal_pptx_with_slides, specs)

    assert set(pngs.keys()) == set(specs)
    for spec, png in pngs.items():
        # PNG signature: 89 50 4E 47 0D 0A 1A 0A
        assert png.startswith(b"\x89PNG\r\n\x1a\n"), (
            f"spec {spec}: bytes don't start with PNG signature: {png[:8]!r}"
        )
        # Sanity size — a real layout render is at least ~3KB, even
        # mostly-empty Office defaults.
        assert len(png) > 1_000, f"spec {spec}: PNG suspiciously small ({len(png)} bytes)"


@sidecar_required
@pytest.mark.asyncio
async def test_render_layouts_handles_empty_specs(minimal_pptx: bytes) -> None:
    """Edge case: caller passes an empty list. We return an empty
    dict without making any HTTP call (saves a roundtrip)."""
    from app.services.pptx_renderer import PptxRenderer

    renderer = PptxRenderer(base_url="http://127.0.0.1:8002")
    pngs = await renderer.render_layouts(minimal_pptx, [])
    assert pngs == {}


@sidecar_required
@pytest.mark.asyncio
async def test_render_layouts_synthesises_unused_layouts(
    minimal_pptx: bytes,
) -> None:
    """Phase 2.4.1: a deck with zero slides exercises the synthesised
    pass — every requested layout is unused-by-source, so the
    renderer must build a synth pptx with one empty slide per spec
    and ship that to the sidecar.
    """
    from app.services.pptx_renderer import PptxRenderer

    renderer = PptxRenderer(base_url="http://127.0.0.1:8002")
    specs = [(0, 0), (0, 1)]
    pngs = await renderer.render_layouts(minimal_pptx, specs)

    # Both specs come back even though the source has no representative
    # slides for them.
    assert set(pngs.keys()) == set(specs)
    for spec, png in pngs.items():
        assert png.startswith(b"\x89PNG\r\n\x1a\n")
        assert len(png) > 1_000

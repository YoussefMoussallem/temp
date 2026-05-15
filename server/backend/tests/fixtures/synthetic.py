"""Synthetic .pptx fixtures built programmatically with python-pptx.

Real corporate templates can't ship in this repo (IP), but tests need
deterministic inputs to be useful. We generate small .pptx files in
memory with controlled properties — known canvas size, theme palette,
layout names, placeholder geometry — so extractor + renderer behaviour
can be asserted exactly.

Each builder returns ``bytes`` (the raw .pptx ZIP) so tests can pass
them directly into ``extract_master_from_pptx`` or save them via
``Presentation.save`` for round-trip checks.

Builders are intentionally minimal — one builder per *behaviour* we
care about (color resolution, layout classification, etc.), not one
per real template. Add a new builder when a new behaviour needs
coverage; resist the urge to make existing builders 'more realistic.'
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Emu


def _new_presentation_16x9() -> Presentation:
    """A blank Presentation sized to standard PowerPoint 16:9 = 1280×720
    CSS px (= 12192000 × 6858000 EMU). python-pptx's default is 4:3
    which is uncommon in modern decks; standardising tests on 16:9 keeps
    geometry assertions intuitive.
    """
    prs = Presentation()
    prs.slide_width = Emu(1280 * 9525)
    prs.slide_height = Emu(720 * 9525)
    return prs


def minimal_master_bytes() -> bytes:
    """The smallest valid .pptx we can produce. Zero slides.

    Useful for tests that need to exercise the *master/layout* path
    without any content slides interfering — extractor + temp-pptx
    builder unit tests.
    """
    prs = _new_presentation_16x9()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def minimal_master_with_slides_bytes() -> bytes:
    """Like ``minimal_master_bytes`` but with one example slide on
    each of the first three layouts.

    Use this when a test needs ``find_layout_representative_slides``
    to actually find something. The integration test for
    ``PptxRenderer`` lives here because the renderer's preview
    strategy is "render the slide that uses this layout"; empty
    decks have no representatives and skip out.
    """
    prs = _new_presentation_16x9()
    layouts = list(prs.slide_masters[0].slide_layouts)
    for li in range(min(3, len(layouts))):
        prs.slides.add_slide(layouts[li])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


__all__ = ["minimal_master_bytes", "minimal_master_with_slides_bytes"]
